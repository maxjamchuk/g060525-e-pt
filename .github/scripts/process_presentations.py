#!/usr/bin/env python3
import sys
import os
import multiprocessing
from concurrent.futures import ProcessPoolExecutor
from functools import partial
from typing import Dict, List, Tuple
from pptx import Presentation

def get_presentation_toc(file_path: str) -> Tuple[str, List[str]]:
    """Получает оглавление презентации."""
    try:
        prs = Presentation(file_path)
        title = file_path.split('/')[-1].replace('.pptx', '')
        slides = []
        for i, slide in enumerate(prs.slides, 1):
            slide_title = ""
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    slide_title = shape.text.strip()
                    if slide_title:
                        break
            if not slide_title:
                slide_title = f"Слайд {i}"
            slides.append(slide_title)
        return title, slides
    except Exception as e:
        return file_path.split('/')[-1].replace('.pptx', ''), [f"Ошибка при обработке: {str(e)}"]

def get_presentation_metadata(file_path: str) -> Tuple[str, Dict]:
    """Получает метаданные презентации."""
    try:
        prs = Presentation(file_path)
        return file_path.split('/')[-1].replace('.pptx', ''), {
            'slides': len(prs.slides),
            'date': prs.core_properties.created.strftime("%Y-%m-%d") if prs.core_properties.created else "Unknown"
        }
    except Exception as e:
        return file_path.split('/')[-1].replace('.pptx', ''), {'slides': 0, 'date': 'Unknown'}

def process_directory(dir_path: str) -> Tuple[str, List[Tuple[str, List[str]]], List[Tuple[str, Dict]]]:
    """Обрабатывает все презентации в директории параллельно."""
    pptx_files = [os.path.join(dir_path, f) for f in os.listdir(dir_path) if f.endswith('.pptx')]
    
    if not pptx_files:
        return dir_path, [], []
    
    # Определяем количество процессов (не более количества файлов)
    num_processes = min(multiprocessing.cpu_count(), len(pptx_files))
    
    with ProcessPoolExecutor(max_workers=num_processes) as executor:
        # Параллельно получаем оглавления
        toc_futures = list(executor.map(get_presentation_toc, pptx_files))
        
        # Параллельно получаем метаданные
        metadata_futures = list(executor.map(get_presentation_metadata, pptx_files))
    
    return dir_path, toc_futures, metadata_futures

def create_presentations_md(dir_path: str, toc_data: List[Tuple[str, List[str]]]) -> None:
    """Создает файл presentations.md в директории."""
    md_file = os.path.join(dir_path, 'presentations.md')
    
    with open(md_file, 'w', encoding='utf-8') as f:
        f.write(f"# Презентации в директории {os.path.basename(dir_path)}\n\n")
        
        for title, slides in toc_data:
            f.write(f"## {title}\n\n")
            for slide in slides:
                f.write(f"- {slide}\n")
            f.write("\n")

def main():
    if len(sys.argv) != 2:
        print("Использование: python process_presentations.py <путь_к_директории>")
        sys.exit(1)
    
    dir_path = sys.argv[1]
    if not os.path.isdir(dir_path):
        print(f"Ошибка: {dir_path} не является директорией")
        sys.exit(1)
    
    # Обрабатываем директорию
    dir_path, toc_data, metadata = process_directory(dir_path)
    
    # Создаем presentations.md
    if toc_data:
        create_presentations_md(dir_path, toc_data)
    
    # Выводим метаданные в формате для bash
    for title, meta in metadata:
        print(f"{title}|{meta['slides']}|{meta['date']}")

if __name__ == "__main__":
    main() 