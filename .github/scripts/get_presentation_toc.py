#!/usr/bin/env python3
import sys
from pptx import Presentation

def get_presentation_toc(file_path):
    try:
        prs = Presentation(file_path)
        print(f"## {file_path.split('/')[-1].replace('.pptx', '')}")
        print("")
        for i, slide in enumerate(prs.slides, 1):
            title = ""
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    title = shape.text.strip()
                    if title:
                        break
            if not title:
                title = f"Слайд {i}"
            print(f"- {title}")
        print("")
    except Exception as e:
        print(f"Ошибка при обработке презентации: {str(e)}")

if __name__ == "__main__":
    if len(sys.argv) != 2:
        print("Использование: python get_presentation_toc.py <путь_к_презентации>")
        sys.exit(1)
    get_presentation_toc(sys.argv[1]) 